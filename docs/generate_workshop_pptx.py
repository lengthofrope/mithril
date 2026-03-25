"""Generate PowerPoint presentation for TDD, ADR & AI-Assisted Development workshop."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
SAGE_GREEN = RGBColor(0x3D, 0x8B, 0x6B)
DARK_GREEN = RGBColor(0x2A, 0x5F, 0x4A)
WARM_STONE = RGBColor(0x4A, 0x45, 0x40)
LIGHT_BG = RGBColor(0xF5, 0xF2, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_AMBER = RGBColor(0xC4, 0x9A, 0x3C)
LIGHT_SAGE = RGBColor(0xE8, 0xF0, 0xEC)
MID_GRAY = RGBColor(0x6B, 0x65, 0x5E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_background(slide, color=LIGHT_BG):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, color=SAGE_GREEN):
    """Add a thin accent bar at the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_side_block(slide, color=SAGE_GREEN):
    """Add a colored block on the left side."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title_text(slide, text, left, top, width, height, size=36, color=WARM_STONE, bold=True):
    """Add a title text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return tf


def add_body_text(slide, text, left, top, width, height, size=20, color=WARM_STONE):
    """Add body text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return tf


def add_bullet_list(slide, items, left, top, width, height, size=20, color=WARM_STONE, spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return tf


def add_section_slide(slide, section_num, section_title, subtitle=""):
    """Create a section divider slide."""
    add_background(slide, SAGE_GREEN)

    # Section number
    add_title_text(slide, f"DEEL {section_num}", Inches(1.5), Inches(2), Inches(10), Inches(1),
                   size=24, color=RGBColor(0xC8, 0xDE, 0xD2), bold=True)

    # Title
    add_title_text(slide, section_title, Inches(1.5), Inches(2.7), Inches(10), Inches(1.5),
                   size=48, color=WHITE, bold=True)

    if subtitle:
        add_body_text(slide, subtitle, Inches(1.5), Inches(4.3), Inches(8), Inches(1),
                      size=22, color=RGBColor(0xC8, 0xDE, 0xD2))


def add_card(slide, left, top, width, height, title, items, title_color=SAGE_GREEN):
    """Add a card-style grouped element."""
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xD8, 0xD3, 0xCC)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    # Title
    add_title_text(slide, title, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.6),
                   size=20, color=title_color, bold=True)

    # Items
    if items:
        add_bullet_list(slide, items, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6),
                        height - Inches(0.9), size=16, color=WARM_STONE, spacing=Pt(4))


# =============================================================================
# SLIDE 1: Title
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, SAGE_GREEN)

add_title_text(slide, "TDD, ADR &\nAI-Assisted Development",
               Inches(1.5), Inches(1.5), Inches(10), Inches(2.5),
               size=52, color=WHITE, bold=True)

add_body_text(slide, "Hoe je voorkomt dat je AI-assistent je codebase sloopt",
              Inches(1.5), Inches(4.2), Inches(8), Inches(0.8),
              size=24, color=RGBColor(0xC8, 0xDE, 0xD2))

# Divider line
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.2), Inches(2), Inches(0.04)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_AMBER
shape.line.fill.background()

add_body_text(slide, "Case study: Mithril\n418 commits  |  1091 tests  |  28 ADRs  |  3 weken",
              Inches(1.5), Inches(5.5), Inches(8), Inches(1),
              size=18, color=RGBColor(0xC8, 0xDE, 0xD2))


# =============================================================================
# SLIDE 2: Agenda
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Agenda", Inches(1), Inches(0.6), Inches(10), Inches(0.8),
               size=40, color=WARM_STONE)

items = [
    ("1", "Het Probleem", "Context Rot"),
    ("2", "De Vier Lagen", "PRD, Plan, TDD, ADR"),
    ("3", "Toepassen", "Lessen en valkuilen"),
    ("4", "ANVIL", "Automatisering van de flow"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.8) + Inches(i * 1.2)

    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAGE_GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_title_text(slide, title, Inches(2.5), y + Inches(0.0), Inches(5), Inches(0.45),
                   size=26, color=WARM_STONE, bold=True)
    add_body_text(slide, desc, Inches(2.5), y + Inches(0.4), Inches(5), Inches(0.35),
                  size=18, color=MID_GRAY)


# =============================================================================
# SLIDE 3: Section - Context Rot
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "1", "Het Probleem", "Context Rot")


# =============================================================================
# SLIDE 4: AI heeft geheugenverlies
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "AI heeft geheugenverlies", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "AI werkt binnen een context window; een beperkt geheugen. Elke chat begint blanco.",
              Inches(1), Inches(1.4), Inches(11), Inches(0.6),
              size=20, color=MID_GRAY)

add_card(slide, Inches(1), Inches(2.3), Inches(3.5), Inches(3.5),
         "Context compressie",
         ["Oudere berichten worden", "samengevat of weggelaten", "Nuances en afspraken", "verdwijnen stilletjes"])

add_card(slide, Inches(5), Inches(2.3), Inches(3.5), Inches(3.5),
         "Attention decay",
         ["Meer tekst = minder gewicht", "voor oudere informatie", "Instructies van 20 berichten", "geleden zijn effectief vergeten"])

add_card(slide, Inches(9), Inches(2.3), Inches(3.5), Inches(3.5),
         "Hallucinaties",
         ["AI vult gaten met", "plausibele maar onjuiste info", "Weet niet dat hij het", "niet meer weet"])

# Quote
tf = add_body_text(slide, "\"Stel je voor dat je een collega hebt die scherp begint, maar gedurende de dag\nsteeds meer vergeet; en het ergste is dat hij niet doorheeft dat hij dingen vergeet.\"",
              Inches(1), Inches(6.2), Inches(11), Inches(0.9),
              size=18, color=MID_GRAY)
tf.paragraphs[0].font.italic = True


# =============================================================================
# SLIDE 5: De gevolgen
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Context Rot in de praktijk", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

problems = [
    "Architectuurbeslissingen gaan verloren of raken verdund",
    "De AI maakt dezelfde fouten opnieuw",
    "Code wordt inconsistent; patronen worden vergeten",
    "\"Waarom is dit zo?\" is niet meer te beantwoorden",
    "AI stelt code voor die ingaat tegen eerdere afspraken",
]
add_bullet_list(slide, problems, Inches(1), Inches(1.8), Inches(10), Inches(4),
                size=24, color=WARM_STONE, spacing=Pt(18))


# =============================================================================
# SLIDE 6: De oplossing - 4 lagen
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "De oplossing: externaliseer context", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "De AI vergeet, maar bestanden in je repository niet.",
              Inches(1), Inches(1.4), Inches(11), Inches(0.6),
              size=22, color=MID_GRAY)

layers = [
    ("PRD", "Wat & Waarom", "Feature-eisen, doelen,\nacceptatiecriteria"),
    ("Plan", "Hoe", "Technische scope,\nfases en voortgang"),
    ("TDD", "Werkt het", "Gedrag, contracten,\nregressiebeveiliging"),
    ("ADR", "Waarom zo", "Architectuurbeslissingen\nen onderbouwing"),
]

for i, (title, question, desc) in enumerate(layers):
    x = Inches(1) + Inches(i * 3.1)
    # Card bg
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(2.8), Inches(3.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xD8, 0xD3, 0xCC)
    shape.line.width = Pt(1)

    # Layer label
    label_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), Inches(2.7), Inches(2.2), Inches(0.55))
    label_shape.fill.solid()
    label_shape.fill.fore_color.rgb = SAGE_GREEN
    label_shape.line.fill.background()
    ltf = label_shape.text_frame
    p = ltf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Question
    add_title_text(slide, f"Beantwoordt: {question}", x + Inches(0.3), Inches(3.5), Inches(2.2), Inches(0.5),
                   size=16, color=ACCENT_AMBER, bold=True)

    add_body_text(slide, desc, x + Inches(0.3), Inches(4.2), Inches(2.2), Inches(1.5),
                  size=17, color=WARM_STONE)


# =============================================================================
# SLIDE 7: Section - De Vier Lagen
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "2", "De Vier Lagen in de Praktijk")


# =============================================================================
# SLIDE 8: Laag 1 - PRDs
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Laag 1: PRDs; de opdracht", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "Wat gaan we bouwen en waarom?",
              Inches(1), Inches(1.4), Inches(11), Inches(0.6),
              size=22, color=MID_GRAY)

prd_items = [
    "Doel; welk probleem lost dit op?",
    "Scope; wat is in en out of scope?",
    "Acceptatiecriteria; wanneer is het klaar?",
]
add_bullet_list(slide, prd_items, Inches(1), Inches(2.5), Inches(5.5), Inches(2.5),
                size=22, color=WARM_STONE, spacing=Pt(16))

# Lesson box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.0), Inches(11), Inches(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_SAGE
shape.line.fill.background()

tf = add_body_text(slide, "Een PRD hoeft niet zwaar te zijn. Zelfs een half A4 met doel, scope en\n5 acceptatiecriteria bespaart uren aan verkeerde implementatie.",
              Inches(1.3), Inches(5.2), Inches(10.4), Inches(1.1),
              size=20, color=DARK_GREEN)
tf.paragraphs[0].font.italic = True


# =============================================================================
# SLIDE 9: Laag 2 - Het Plan
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Laag 2: Het Plan; de vertaling", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "Vertaalt PRD-criteria naar technische specs",
              Inches(1), Inches(1.4), Inches(11), Inches(0.6),
              size=22, color=MID_GRAY)

plan_items = [
    "Deelt werk op in fases die in een sessie passen",
    "Maakt elke spec testbaar (directe input voor TDD)",
    "Houdt voortgang bij; legt vast wat out of scope is",
    "Levend document; wordt bijgewerkt tijdens implementatie",
]
add_bullet_list(slide, plan_items, Inches(1), Inches(2.5), Inches(10), Inches(3),
                size=22, color=WARM_STONE, spacing=Pt(16))

# Example
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.0), Inches(11), Inches(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_SAGE
shape.line.fill.background()

tf = add_body_text(slide, "Voorbeeld: TDD ontdekte een event dispatch gap; plan werd bijgewerkt\nmet een nieuwe Phase 0, en ADR-013 geschreven.",
              Inches(1.3), Inches(5.2), Inches(10.4), Inches(1.1),
              size=20, color=DARK_GREEN)
tf.paragraphs[0].font.italic = True


# =============================================================================
# SLIDE 10: Laag 3 - TDD
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Laag 3: TDD; het veiligheidsnet", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "Geen \"nice to have\"; het veiligheidsnet dat voorkomt dat de AI je codebase sloopt.",
              Inches(1), Inches(1.4), Inches(11), Inches(0.6),
              size=22, color=MID_GRAY)

# Red-Green-Refactor cards
colors_rgb = [RGBColor(0xC0, 0x39, 0x2B), SAGE_GREEN, RGBColor(0x29, 0x80, 0xB9)]
titles = ["RED", "GREEN", "REFACTOR"]
descs = [
    "Schrijf een falende test\ndie het gewenste gedrag\nbeschrijft",
    "Schrijf de minimale code\nom de test te laten\nslagen",
    "Verbeter de code met\nhet vangnet van\ngroene tests",
]

for i in range(3):
    x = Inches(1) + Inches(i * 4.1)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(3.6), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = colors_rgb[i]
    shape.line.width = Pt(2)

    # Colored header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.5), Inches(3.6), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors_rgb[i]
    bar.line.fill.background()
    btf = bar.text_frame
    p = btf.paragraphs[0]
    p.text = titles[i]
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_body_text(slide, descs[i], x + Inches(0.3), Inches(3.5), Inches(3), Inches(1.8),
                  size=18, color=WARM_STONE)

# Arrow connectors between cards
for i in range(2):
    x = Inches(4.4) + Inches(i * 4.1)
    add_title_text(slide, ">", x, Inches(3.3), Inches(0.5), Inches(0.5),
                   size=36, color=MID_GRAY, bold=True)


# =============================================================================
# SLIDE 11: TDD Case Study
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Case Study: HasFollowUp trait (ADR-001)", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=34, color=WARM_STONE)

# Problem card
add_card(slide, Inches(1), Inches(1.6), Inches(5.3), Inches(3),
         "Het probleem",
         ["HasFollowUp trait bevroeg kolommen", "follow_up_date en status direct",
          "Die kolommen bestaan alleen op de", "follow_ups tabel; niet op tasks",
          "Runtime SQL errors gegarandeerd"],
         title_color=RGBColor(0xC0, 0x39, 0x2B))

# Solution card
add_card(slide, Inches(7), Inches(1.6), Inches(5.3), Inches(3),
         "De oplossing",
         ["Test faalde: niet de implementatie", "was fout, maar het hele design",
          "Scopes herschreven naar", "whereHas-based queries",
          "Bug gevonden voor productcode"],
         title_color=SAGE_GREEN)

# Lesson
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.2), Inches(11), Inches(1.3))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_SAGE
shape.line.fill.background()

tf = add_body_text(slide, "De tests vonden de bug. Niet de AI, niet code review, niet QA.\nDe falende test dwong een redesign af.",
              Inches(1.3), Inches(5.4), Inches(10.4), Inches(1),
              size=22, color=DARK_GREEN)
tf.paragraphs[0].font.bold = True


# =============================================================================
# SLIDE 12: Laag 4 - ADRs
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Laag 4: ADRs; het geheugen", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "Documenteert beslissingen die afwijken van het plan of niet-triviale trade-offs bevatten.",
              Inches(1), Inches(1.4), Inches(11), Inches(0.6),
              size=22, color=MID_GRAY)

# Three pillars
adr_titles = ["Context", "Decision", "Consequences"]
adr_descs = ["Wat was het probleem?", "Wat is er besloten\nen waarom?", "Wat verandert er?"]

for i in range(3):
    x = Inches(1) + Inches(i * 4.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(3.6), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xD8, 0xD3, 0xCC)
    shape.line.width = Pt(1)

    add_title_text(slide, adr_titles[i], x + Inches(0.3), Inches(2.7), Inches(3), Inches(0.5),
                   size=24, color=SAGE_GREEN, bold=True)
    add_body_text(slide, adr_descs[i], x + Inches(0.3), Inches(3.3), Inches(3), Inches(0.8),
                  size=20, color=WARM_STONE)

# Key insight
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.2), Inches(11), Inches(1.3))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_SAGE
shape.line.fill.background()

tf = add_body_text(slide, "Vuistregel: Als iemand later zou kunnen vragen \"waarom is dit zo?\"; schrijf een ADR.",
              Inches(1.3), Inches(5.4), Inches(10.4), Inches(0.9),
              size=22, color=DARK_GREEN)
tf.paragraphs[0].font.bold = True


# =============================================================================
# SLIDE 13: ADR Case Study
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Case Study: Event Dispatch Gap (ADR-013)", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=34, color=WARM_STONE)

add_card(slide, Inches(1), Inches(1.6), Inches(5.3), Inches(2.5),
         "Het probleem",
         ["TaskStatusChanged event alleen", "gefired vanuit expliciete dispatch",
          "3 andere code paths updaten status", "zonder event te dispatchen"],
         title_color=RGBColor(0xC0, 0x39, 0x2B))

add_card(slide, Inches(7), Inches(1.6), Inches(5.3), Inches(2.5),
         "De oplossing",
         ["TaskObserver dispatcht event", "automatisch bij elke statuswijziging",
          "Handmatige dispatches verwijderd", "Loste ook pre-bestaand probleem op"],
         title_color=SAGE_GREEN)

# What the ADR captures
add_body_text(slide, "Zonder ADR zou een volgende sessie niet weten:",
              Inches(1), Inches(4.5), Inches(11), Inches(0.5),
              size=20, color=WARM_STONE)

adr_items = [
    "Dat er bewust een Observer is gekozen boven handmatige dispatch calls",
    "Dat handmatige dispatches bewust zijn verwijderd",
    "Welke alternatieven zijn afgewezen en waarom",
]
add_bullet_list(slide, adr_items, Inches(1), Inches(5.1), Inches(11), Inches(2),
                size=18, color=WARM_STONE, spacing=Pt(8))


# =============================================================================
# SLIDE 14: De Feedback Loop
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "De Feedback Loop", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

# Flow diagram using shapes
flow_steps = [
    "PRD definieert\nwat en waarom",
    "Plan vertaalt naar\ntechnische specs",
    "Specs worden\ntests (TDD)",
    "Tests ontdekken\nproblemen",
    "Problemen\ntriggeren ADRs",
    "ADRs informeren\nvolgende sessies",
]

for i, step in enumerate(flow_steps):
    x = Inches(0.7) + Inches(i * 2.1)
    y = Inches(2)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.85), Inches(1.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAGE_GREEN if i % 2 == 0 else WHITE
    shape.line.color.rgb = SAGE_GREEN
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE if i % 2 == 0 else SAGE_GREEN
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if i < 5:
        add_title_text(slide, ">", x + Inches(1.85), y + Inches(0.35), Inches(0.3), Inches(0.5),
                       size=28, color=SAGE_GREEN, bold=True)

# Problem-solution table
headers = [("Probleem", "Oplossing")]
rows = [
    ("AI vergeet vorige sessie", "ADRs + plan overleven sessies"),
    ("AI breekt bestaande code", "Tests vangen regressies direct"),
    ("AI gebruikt andere patronen", "Plan + codebase scan dwingen af"),
    ("AI voegt ongeplande features toe", "Plan met \"Out of Scope\" sectie"),
    ("AI draait bewuste keuzes terug", "ADRs documenteren het waarom"),
]

table_top = Inches(4)
table = slide.shapes.add_table(len(rows) + 1, 2, Inches(1.5), table_top, Inches(10), Inches(3)).table
table.columns[0].width = Inches(5)
table.columns[1].width = Inches(5)

# Header
for j, (h1, h2) in enumerate(headers):
    for col_idx, text in enumerate([h1, h2]):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = SAGE_GREEN
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True

# Rows
for i, (prob, sol) in enumerate(rows):
    for col_idx, text in enumerate([prob, sol]):
        cell = table.cell(i + 1, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_SAGE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(15)
        p.font.color.rgb = WARM_STONE


# =============================================================================
# SLIDE 15: Section - Toepassen
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "3", "Toepassen en Valkuilen")


# =============================================================================
# SLIDE 16: De vijf lessen
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "De vijf belangrijkste lessen", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

lessons = [
    ("1", "Als het niet in een bestand staat, bestaat het niet",
     "Elke AI-sessie begint leeg. Context die alleen in je hoofd zit, gaat verloren."),
    ("2", "Maak specs testbaar",
     "Als je het niet in een test kunt uitdrukken, is de spec te vaag."),
    ("3", "Schrijf ADRs op het moment van de beslissing",
     "Achteraf mis je de alternatieven die je hebt overwogen."),
    ("4", "Houd sessies klein",
     "Kleine fases houden sessies beheersbaar en voorkomen context compressie."),
    ("5", "Jij blijft de architect",
     "Review, denk mee, stel vragen. Blind vertrouwen leidt tot blinde vlekken."),
]

for i, (num, title, desc) in enumerate(lessons):
    y = Inches(1.5) + Inches(i * 1.15)

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y + Inches(0.05), Inches(0.55), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAGE_GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_title_text(slide, title, Inches(1.8), y, Inches(10), Inches(0.45),
                   size=22, color=WARM_STONE, bold=True)
    add_body_text(slide, desc, Inches(1.8), y + Inches(0.45), Inches(10), Inches(0.4),
                  size=17, color=MID_GRAY)


# =============================================================================
# SLIDE 17: Hoe begin je?
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Hoe begin je?", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "Je hebt geen speciale tooling nodig",
              Inches(1), Inches(1.4), Inches(11), Inches(0.6),
              size=22, color=MID_GRAY)

how_items = [
    "Start elke AI-sessie met context laden; geef relevante ADRs, het plan, en laat tests lezen",
    "Gebruik een CLAUDE.md (of equivalent); conventies, actieve plannen, verwijzingen naar ADRs",
    "Houd ADRs klein; een ADR per beslissing, niet een mega-document per feature",
    "Gebruik het plan als communicatiemiddel; niet alleen voor de AI, maar ook voor je team",
]
add_bullet_list(slide, how_items, Inches(1), Inches(2.5), Inches(11), Inches(4),
                size=22, color=WARM_STONE, spacing=Pt(20))


# =============================================================================
# SLIDE 18: Mithril in cijfers
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Mithril in cijfers", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "Gebouwd met consistente architectuur ondanks tientallen sessiewisselingen",
              Inches(1), Inches(1.3), Inches(11), Inches(0.6),
              size=20, color=MID_GRAY)

metrics = [
    ("418", "Commits"),
    ("187", "Testbestanden"),
    ("1091", "Tests"),
    ("28", "ADRs"),
    ("19", "Plannen"),
    ("3", "Weken"),
]

for i, (number, label) in enumerate(metrics):
    x = Inches(1) + Inches(i * 2.05)
    y = Inches(2.5)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.8), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xD8, 0xD3, 0xCC)
    shape.line.width = Pt(1)

    add_title_text(slide, number, x + Inches(0.1), y + Inches(0.4), Inches(1.6), Inches(1),
                   size=44, color=SAGE_GREEN, bold=True)

    tf = add_body_text(slide, label, x + Inches(0.1), y + Inches(1.5), Inches(1.6), Inches(0.6),
                       size=18, color=WARM_STONE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Center the number too
    # Re-get the number text frame to center
    for shp in slide.shapes:
        pass  # already added inline


# =============================================================================
# SLIDE 19: Section - ANVIL
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, "4", "ANVIL", "Automated Nexus for Verified Iterative Lifecycles")


# =============================================================================
# SLIDE 20: Van handmatig naar georchestreerd
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Van handmatig naar georchestreerd", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "De vier lagen werken, maar vereisen discipline. Onder tijdsdruk slipt dat.",
              Inches(1), Inches(1.4), Inches(11), Inches(0.6),
              size=22, color=MID_GRAY)

# ANVIL flow
anvil_steps = [
    ("Research", "Analyst"),
    ("PRD", "Specifier"),
    ("Plan", "Architect"),
    ("Build", "Builder (TDD)"),
    ("ADRs", "Automatisch"),
]

for i, (step, role) in enumerate(anvil_steps):
    x = Inches(0.8) + Inches(i * 2.5)
    y = Inches(2.8)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.1), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAGE_GREEN
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = role
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0xC8, 0xDE, 0xD2)
    p2.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if i < 4:
        add_title_text(slide, ">", x + Inches(2.1), y + Inches(0.4), Inches(0.4), Inches(0.5),
                       size=28, color=SAGE_GREEN, bold=True)

# Additional info
add_body_text(slide, "Elke rol is een AI-agent met specifieke opdracht en regels.\nDe orchestrator stuurt aan, bewaakt voortgang, en houdt de gebruiker in controle.",
              Inches(1), Inches(5), Inches(11), Inches(1),
              size=20, color=WARM_STONE)


# =============================================================================
# SLIDE 21: ANVIL Rollen
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Zes gespecialiseerde rollen", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

roles = [
    ("Analyst", "Haalt requirements op uit\nexterne bronnen (bijv. Jira)"),
    ("Specifier", "Schrijft en reviewt PRDs"),
    ("Architect", "Maakt implementatieplannen\nmet fases en specs"),
    ("Builder", "Bouwt via TDD, kan parallel\nwerken in agent teams"),
    ("Auditor", "OWASP security scans\nop de codebase"),
    ("Scribe", "Genereert documenten vanuit\nde opgeleverde resultaten"),
]

for i, (role, desc) in enumerate(roles):
    col = i % 3
    row = i // 3
    x = Inches(1) + Inches(col * 4)
    y = Inches(1.6) + Inches(row * 2.8)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.5), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xD8, 0xD3, 0xCC)
    shape.line.width = Pt(1)

    add_title_text(slide, role, x + Inches(0.3), y + Inches(0.2), Inches(3), Inches(0.5),
                   size=24, color=SAGE_GREEN, bold=True)
    add_body_text(slide, desc, x + Inches(0.3), y + Inches(0.8), Inches(3), Inches(1.2),
                  size=17, color=WARM_STONE)


# =============================================================================
# SLIDE 22: Automatische bewaking
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Automatische bewaking via Hooks", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

add_body_text(slide, "Twee PostToolUse hooks die real-time afdwingen wat anders op discipline leunt",
              Inches(1), Inches(1.4), Inches(11), Inches(0.6),
              size=20, color=MID_GRAY)

# TDD Enforce card
add_card(slide, Inches(1), Inches(2.3), Inches(5.3), Inches(3.2),
         "TDD Enforce",
         ["Waarschuwt wanneer bronbestand", "gewijzigd wordt zonder testwijzigingen",
          "Dwingt Red-Green volgorde af:", "eerst tests, dan implementatie",
          "Melding op moment van edit,", "niet pas bij commit"],
         title_color=RGBColor(0xC0, 0x39, 0x2B))

# ADR Watch card
add_card(slide, Inches(7), Inches(2.3), Inches(5.3), Inches(3.2),
         "ADR Watch",
         ["Waarschuwt wanneer bestand", "buiten actief plan wordt gewijzigd",
          "Voorkomt onbewuste scope creep",
          "Informatief; blokkeert niet"],
         title_color=ACCENT_AMBER)

# Note
tf = add_body_text(slide, "Beide hooks zijn informatief (ze blokkeren niet) en worden automatisch geinstalleerd.\nDit zijn mechanische checks die onafhankelijk van de AI draaien.",
              Inches(1), Inches(6), Inches(11), Inches(1),
              size=18, color=MID_GRAY)
tf.paragraphs[0].font.italic = True


# =============================================================================
# SLIDE 23: Context rot oplossingen
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "Hoe ANVIL context rot minimaliseert", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=36, color=WARM_STONE)

# Three mechanism cards
mechanisms = [
    ("State op disk", "Lifecycle state afgeleid uit\nbestanden, niet uit geheugen.\nSession handoff bij einde sessie."),
    ("Kleine, gefocuste agents", "Elke agent krijgt alleen wat\nhij nodig heeft. Kort window =\nminder attention decay."),
    ("Mechanische bewaking", "Hooks vangen fouten door\ncontext rot op het moment\ndat ze gebeuren."),
]

for i, (title, desc) in enumerate(mechanisms):
    x = Inches(1) + Inches(i * 4.1)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.6), Inches(2.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = SAGE_GREEN
    shape.line.width = Pt(1.5)

    # Number
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), Inches(2), Inches(0.5), Inches(0.5))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = SAGE_GREEN
    num_shape.line.fill.background()
    ntf = num_shape.text_frame
    p = ntf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_title_text(slide, title, x + Inches(0.8), Inches(2), Inches(2.5), Inches(0.5),
                   size=20, color=WARM_STONE, bold=True)
    add_body_text(slide, desc, x + Inches(0.3), Inches(2.7), Inches(3), Inches(1.5),
                  size=16, color=WARM_STONE)

# Summary table
table = slide.shapes.add_table(6, 2, Inches(1), Inches(5), Inches(11), Inches(2.3)).table
table.columns[0].width = Inches(5.5)
table.columns[1].width = Inches(5.5)

table_data = [
    ("Context rot probleem", "ANVIL oplossing"),
    ("Sessiewissel verliest voortgang", "Session handoff + state op disk"),
    ("Lang window; attention decay", "Kleine agents met gefocuste context"),
    ("AI vergeet eigen regels", "Hooks onafhankelijk van de AI"),
    ("AI herhaalt afgewezen keuzes", "ADRs geladen in elke agent prompt"),
    ("Scope creep", "ADR Watch detecteert out-of-plan edits"),
]

for i, (col1, col2) in enumerate(table_data):
    for j, text in enumerate([col1, col2]):
        cell = table.cell(i, j)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SAGE_GREEN
            p.font.color.rgb = WHITE
            p.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LIGHT_SAGE
            p.font.color.rgb = WARM_STONE


# =============================================================================
# SLIDE 24: De kern blijft
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_accent_bar(slide)

add_title_text(slide, "De kern blijft", Inches(1), Inches(0.5), Inches(11), Inches(0.8),
               size=38, color=WARM_STONE)

# Big quote style
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2), Inches(10), Inches(3))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = SAGE_GREEN
shape.line.width = Pt(2)

tf = add_body_text(slide, "ANVIL automatiseert de flow, maar de principes zijn dezelfde:\nPRDs, plannen, tests en beslissingen; vastgelegd in bestanden\ndie elke sessie overleven.",
              Inches(2), Inches(2.5), Inches(9), Inches(1.5),
              size=26, color=WARM_STONE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

tf2 = add_body_text(slide, "De tooling maakt het makkelijker om de discipline vol te houden,\nmaar je kunt dezelfde aanpak toepassen met elk AI-hulpmiddel\nen een teksteditor.",
              Inches(2), Inches(3.8), Inches(9), Inches(1.2),
              size=20, color=MID_GRAY)
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER


# =============================================================================
# SLIDE 25: Discussie / Q&A
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, SAGE_GREEN)

add_title_text(slide, "Discussie / Q&A", Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
               size=48, color=WHITE, bold=True)

# Divider
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3), Inches(2), Inches(0.04)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_AMBER
shape.line.fill.background()

qa_items = [
    "Hoe past dit in jullie huidige workflow?",
    "Welke onderdelen zijn direct toepasbaar; ook zonder AI?",
    "Wat is de overhead vs. de tijdsbesparing?",
]
tf = add_bullet_list(slide, qa_items, Inches(1.5), Inches(3.5), Inches(9), Inches(3),
                     size=24, color=RGBColor(0xC8, 0xDE, 0xD2), spacing=Pt(20))


# =============================================================================
# Save
# =============================================================================
output_path = "/home/bdekort/Sites/mithril/docs/workshop-tdd-adr-ai-assisted-development.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
