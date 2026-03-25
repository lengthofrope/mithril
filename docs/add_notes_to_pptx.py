"""Add speaker notes to the workshop PowerPoint based on the markdown content."""

from pptx import Presentation

prs = Presentation("workshop-tdd-adr-ai-assisted-development.pptx")

# Speaker notes per slide (0-indexed)
notes = {
    0: """Welkom bij de workshop TDD, ADR & AI-Assisted Development.

Case study: Mithril; een PWA gebouwd met Claude Code. 418 commits, 1091 tests, 28 ADRs in 3 weken. Dit is een echt project, geen demo.

De workshop duurt ongeveer 35 minuten. Doelgroep: developers die werken met of gaan werken met AI coding assistants.""",

    1: """Korte agenda-overview:
- Deel 1: Het probleem; context rot. Waarom AI-assistants zonder structuur je codebase beschadigen.
- Deel 2: De vier lagen; PRD, Plan, TDD, ADR. Concrete tools om context te externaliseren.
- Deel 3: Toepassen; de vijf belangrijkste lessen en hoe je morgen kunt beginnen.
- Deel 4: ANVIL; hoe je deze flow kunt automatiseren met een orchestrator.""",

    2: """Transitieslide naar Deel 1: Het Probleem; Context Rot.

We beginnen met begrijpen waarom AI-assistants falen bij langere of meerdere sessies.""",

    3: """AI coding assistants werken binnen een context window; een beperkt geheugen met een vaste maximale grootte. Elke nieuwe chat begint blanco.

Maar het probleem gaat verder dan alleen sessiewisselingen. Ook binnen een lopend gesprek degradeert de context op drie manieren:

1. Context compressie: Wanneer een gesprek het window nadert, moet de AI oudere berichten samenvatten of laten vallen. Nuances, constraints en eerdere afspraken verdwijnen stilletjes.

2. Attention decay: Hoe meer tekst er in het window zit, hoe minder gewicht de AI geeft aan informatie die "ver weg" staat. Instructies van 20 berichten geleden worden effectief vergeten, ook al zijn ze technisch nog aanwezig.

3. Hallucinaties: Wanneer de relevante context verdwenen of verdund is, vult de AI de gaten met plausibel klinkende maar feitelijk onjuiste informatie. De AI weet niet dat hij het niet meer weet; hij genereert zelfverzekerd antwoorden op basis van incomplete context.

Vergelijking: Stel je voor dat je een collega hebt die scherp begint, maar gedurende de dag steeds meer vergeet; en het ergste is dat hij niet doorheeft dat hij dingen vergeet.""",

    4: """Het resultaat van deze drie problemen is context rot: een geleidelijk verval van de kwaliteit van AI-output.

In de praktijk betekent dit:
- Architectuurbeslissingen gaan verloren (tussen sessies) of raken verdund (binnen een sessie)
- De AI maakt dezelfde fouten opnieuw
- Code wordt inconsistent omdat patronen vergeten worden
- "Waarom is dit zo?" is niet meer te beantwoorden; niet door jou, niet door de AI
- Na een lange sessie kan de AI zelfverzekerd code voorstellen die indruist tegen afspraken van eerder in datzelfde gesprek""",

    5: """De oplossing: de AI vergeet, maar bestanden in je repository niet. De sleutel is context externaliseren in vier lagen:

- PRDs beantwoorden Wat en Waarom: feature-eisen, doelen en acceptatiecriteria
- Plans beantwoorden Hoe: technische scope, fases en voortgang
- TDD beantwoordt Werkt het: gedrag, contracten en regressiebeveiliging
- ADRs beantwoorden Waarom zo: architectuurbeslissingen en hun onderbouwing

Dit is de kern van de workshop: deze vier lagen gaan we nu stuk voor stuk doorlopen.""",

    6: """Transitieslide naar Deel 2: De Vier Lagen in de Praktijk.

Nu gaan we elke laag bekijken met concrete voorbeelden uit het Mithril project.""",

    7: """Laag 1: PRDs; de opdracht.

Alles begint met de vraag: wat gaan we bouwen en waarom? Een PRD formaliseert dat:
- Doel; welk probleem lost dit op?
- Scope; wat is in en out of scope?
- Acceptatiecriteria; wanneer is het klaar?

Zonder PRD begint de AI (en soms de developer) te bouwen op basis van aannames. De PRD is het contract tussen het idee en de implementatie; het voorkomt dat er een technisch plan wordt geschreven voor iets dat nog niet goed gedefinieerd is.

Belangrijke les: Een PRD hoeft niet zwaar te zijn. Zelfs een half A4 met doel, scope en 5 acceptatiecriteria bespaart uren aan verkeerde implementatie.""",

    8: """Laag 2: Het Plan; de vertaling.

Met een PRD in de hand is de volgende stap: hoe gaan we dit bouwen? Het plan vertaalt PRD-criteria naar technische specs:
- Deelt werk op in fases die in een sessie passen
- Maakt elke spec testbaar (directe input voor TDD)
- Houdt voortgang bij en legt vast wat out of scope is

Het plan is ook een levend document. Voorbeeld: tijdens de implementatie van recurring tasks ontdekte TDD een event dispatch gap die niet in het oorspronkelijke plan stond. Het plan werd bijgewerkt met een nieuwe Phase 0, en een ADR geschreven (ADR-013). Dat voorbeeld komt straks terug.""",

    9: """Laag 3: TDD; het veiligheidsnet.

TDD is geen "nice to have" bij AI-assisted development; het is het veiligheidsnet dat voorkomt dat de AI je codebase sloopt.

Zonder TDD merk je niet dat de AI in sessie 5 iets heeft gebroken dat in sessie 2 werkte. Met TDD worden regressies direct gevangen, en de test suite vormt een levende specificatie die elke sessie overleeft.

Red-Green-Refactor:
1. RED; schrijf een falende test die het gewenste gedrag beschrijft
2. GREEN; schrijf de minimale code om de test te laten slagen
3. REFACTOR; verbeter de code met het vangnet van groene tests""",

    10: """Case study: HasFollowUp trait (ADR-001).

De HasFollowUp trait definieerde scopes die direct kolommen follow_up_date en status bevroegen. Maar die kolommen bestaan alleen op de follow_ups tabel; niet op tasks of team_members die de trait gebruiken. Runtime SQL errors gegarandeerd.

TDD ving dit: de test faalde niet omdat de implementatie fout was, maar omdat het hele design fout was. Scopes werden herschreven naar whereHas-based queries. De bug was gevonden voordat er ook maar een regel productcode was geschreven.

Belangrijke les: De tests vonden de bug. Niet de AI, niet code review, niet QA. De falende test dwong een redesign af.""",

    11: """Laag 4: ADRs; het geheugen.

Tijdens het bouwen worden beslissingen genomen die afwijken van het plan, of die niet-triviale trade-offs bevatten. Een ADR documenteert drie dingen:
1. Context; wat was het probleem?
2. Decision; wat is er besloten en waarom?
3. Consequences; wat verandert er?

Zonder ADRs kan de AI in een volgende sessie een beslissing terugdraaien die je bewust hebt genomen, of een alternatief voorstellen dat je al hebt afgewezen.

Vuistregel: Als iemand later zou kunnen vragen "waarom is dit zo?"; schrijf een ADR.""",

    12: """Case study: Event Dispatch Gap (ADR-013).

Bij het implementeren van recurring tasks bleek dat het TaskStatusChanged event alleen werd gefired vanuit expliciete dispatch calls. Maar drie andere code paths (AutoSaveController, bulkUpdate(), move()) updaten status via $task->update() zonder het event te dispatchen.

De oplossing: een TaskObserver die het event automatisch dispatcht bij elke statuswijziging. Handmatige dispatches bewust verwijderd.

Zonder ADR zou een volgende sessie niet weten:
- Dat er bewust een Observer is gekozen boven handmatige dispatch calls (en waarom)
- Dat handmatige dispatches bewust zijn verwijderd
- Dit ook een pre-bestaand probleem met een andere listener oploste

De ADR legt niet alleen vast wat er is besloten, maar ook welke alternatieven zijn afgewezen en waarom. Dat voorkomt dat de AI (of een developer) dezelfde discussie opnieuw voert.""",

    13: """De feedback loop; dit is hoe de vier lagen samenwerken:

PRD definieert wat en waarom
-> Plan vertaalt naar technische specs
-> Specs worden tests (TDD)
-> Tests ontdekken problemen
-> Problemen triggeren ADRs
-> ADRs informeren volgende sessies
-> Consistente implementatie

Samenvatting van probleem vs. oplossing:
- AI vergeet vorige sessie -> ADRs + plan overleven sessies
- AI breekt bestaande code -> Tests vangen regressies direct
- AI gebruikt andere patronen -> Plan + codebase scan dwingen patronen af
- AI voegt ongeplande features toe -> Plan met "Out of Scope" sectie
- AI draait bewuste keuzes terug -> ADRs documenteren het waarom""",

    14: """Transitieslide naar Deel 3: Toepassen en Valkuilen.""",

    15: """De vijf belangrijkste lessen:

1. Als het niet in een bestand staat, bestaat het niet. Elke AI-sessie begint leeg. Context die alleen in je hoofd zit, gaat verloren.

2. Maak specs testbaar. Als je het niet in een test kunt uitdrukken, is de spec te vaag. Dit is waar plan en TDD samenkomen.

3. Schrijf ADRs op het moment van de beslissing. Achteraf mis je de alternatieven die je hebt overwogen.

4. Houd sessies klein. Context windows hebben limieten. Kleine fases in het plan houden sessies beheersbaar en voorkomen context compressie.

5. Jij blijft de architect. De AI is een krachtig hulpmiddel, maar review, denk mee, en stel vragen. Blind vertrouwen leidt tot blinde vlekken.""",

    16: """Hoe begin je? Je hebt geen speciale tooling nodig om dit toe te passen:

- Start elke AI-sessie met context laden; geef de AI bij de start de relevante ADRs, het plan, en laat hem de tests lezen.
- Gebruik een CLAUDE.md (of equivalent); een bestand in je repo dat de AI automatisch leest bij elke sessie. Hierin staan conventies, actieve plannen, en verwijzingen naar ADRs.
- Houd ADRs klein; een ADR per beslissing, niet een mega-document per feature.
- Gebruik het plan als communicatiemiddel; niet alleen voor de AI, maar ook voor je team.""",

    17: """Mithril in cijfers:
- 418 commits
- 187 testbestanden
- 1091 testfuncties
- 28 ADRs
- 19 implementatieplannen
- 3 weken tijdspanne

Gebouwd met consistente architectuur ondanks tientallen sessiewisselingen. Dit illustreert dat de aanpak schaalbaar is.""",

    18: """Transitieslide naar Deel 4: ANVIL; Automated Nexus for Verified Iterative Lifecycles.

Van handmatig naar geautomatiseerd.""",

    19: """De vier lagen uit deel 2 werken, maar vereisen discipline: je moet zelf onthouden een PRD te schrijven, een plan te maken, TDD te volgen, en ADRs bij te houden. In de praktijk slipt dat; vooral onder tijdsdruk.

ANVIL is een skill voor Claude Code die deze flow automatiseert. Het is in actieve ontwikkeling, gebouwd vanuit de lessen van het Mithril project.

ANVIL is een orchestrator die de volledige lifecycle aanstuurt via gespecialiseerde rollen:
Research (Analyst) -> PRD (Specifier) -> Plan (Architect) -> Implementatie (Builder; TDD verplicht) -> ADRs (automatisch bij afwijkingen)

Elke rol is een AI-agent met een specifieke opdracht en regels. De orchestrator stuurt ze aan, bewaakt de voortgang, en zorgt dat de gebruiker altijd de controle houdt.""",

    20: """Zes gespecialiseerde rollen:

- Analyst: Haalt requirements op uit externe bronnen (bijv. Jira)
- Specifier: Schrijft en reviewt PRDs
- Architect: Maakt implementatieplannen met fases en specs
- Builder: Bouwt via TDD (red-green-refactor), kan parallel werken in agent teams
- Auditor: OWASP security scans op de codebase
- Scribe: Genereert documenten vanuit de opgeleverde resultaten""",

    21: """Automatische bewaking via twee PostToolUse hooks die real-time afdwingen wat anders op discipline leunt:

TDD Enforce: waarschuwt direct wanneer je een bronbestand wijzigt zonder dat er testwijzigingen in de working tree staan. Je krijgt de melding op het moment van de edit, niet pas bij een commit. Dit dwingt de Red-Green volgorde af: eerst tests, dan implementatie.

ADR Watch: waarschuwt wanneer een bestand wordt gewijzigd dat buiten het actieve plan valt. Voorkomt onbewuste scope creep.

Beide hooks zijn informatief (ze blokkeren niet) en worden automatisch geinstalleerd. Dit zijn mechanische checks die onafhankelijk van de AI draaien; geen herinneringen die de AI kan vergeten.""",

    22: """Hoe ANVIL context rot minimaliseert; drie mechanismen:

1. State op disk, niet in het geheugen: ANVIL leidt de lifecycle state af uit bestanden. Niets hangt af van wat de AI "onthoudt". Bij sessie-einde schrijft ANVIL een session handoff. De volgende sessie leest dat bestand en pikt het werk op.

2. Kleine, gefocuste agents: In plaats van een enkele AI-sessie die steeds langer wordt, splitst ANVIL werk op. Een Builder-agent krijgt alleen de specs voor zijn fase, de relevante ADRs, en de projectconventies. Kort, gefocust window = minder attention decay. Bijkomend voordeel: minder tokens per aanroep, lagere kosten, snellere respons.

Kanttekening: de orchestrator zelf accumuleert wel state; bij complexe plannen kan hij tegen dezelfde limieten aanlopen.

3. Mechanische bewaking vervangt geheugen: De hooks vangen fouten door context rot op het moment dat ze gebeuren, onafhankelijk van de AI.""",

    23: """De kern blijft: ANVIL automatiseert de flow, maar de principes zijn dezelfde als in deel 2; PRDs, plannen, tests en beslissingen vastgelegd in bestanden die elke sessie overleven.

De tooling maakt het makkelijker om de discipline vol te houden, maar je kunt dezelfde aanpak toepassen met elk AI-hulpmiddel en een teksteditor.

ANVIL is op dit moment nog in actieve ontwikkeling. Ik zal de skill binnenkort beschikbaar stellen voor iedereen binnen Proud Nerds die hem wil uitproberen. Het is een Claude Code skill; je hebt een werkende Claude Code setup nodig.""",

    24: """Discussie en Q&A. Mogelijke gesprekspunten:

- Hoe past dit in jullie huidige workflow?
- Welke onderdelen zijn direct toepasbaar; ook zonder AI?
- Wat is de overhead vs. de tijdsbesparing?

Laat het gesprek open en informeel. De vier lagen (PRD, Plan, TDD, ADR) zijn ook zonder AI-assistants waardevol; het zijn gewoon goede software engineering practices.""",
}

for idx, notes_text in notes.items():
    slide = prs.slides[idx]
    if not slide.has_notes_slide:
        slide.notes_slide  # Access creates it
    slide.notes_slide.notes_text_frame.text = notes_text

prs.save("workshop-tdd-adr-ai-assisted-development.pptx")
print(f"Done: added notes to {len(notes)} slides.")
