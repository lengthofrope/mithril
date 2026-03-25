"""
Update workshop PowerPoint:
1. Fix test/assertion numbers (2142 tests, 4602 assertions)
2. Add ADR-001 and ADR-013 full-text slides after their case studies
3. Rewrite all speaker notes for read-aloud delivery
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = "workshop-tdd-adr-ai-assisted-development.pptx"
prs = Presentation(PPTX_PATH)

# Colors from existing slides
COLOR_DARK = RGBColor(0x4A, 0x45, 0x40)
COLOR_MUTED = RGBColor(0x6B, 0x65, 0x5E)
COLOR_GREEN = RGBColor(0x3D, 0x8B, 0x6B)


# ============================================================
# STEP 1: Fix numbers on slide 17 ("Mithril in cijfers")
# ============================================================
slide17 = prs.slides[17]
for shape in slide17.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == "1091":
                    run.text = "2142"
                elif run.text.strip() == "187":
                    run.text = "4602"
                elif run.text.strip() == "Testbestanden":
                    run.text = "Assertions"

print("Step 1 done: fixed numbers on slide 17")


# ============================================================
# STEP 2: Create ADR content slides
# ============================================================

def add_adr_slide(prs, title_text, sections):
    """Add a slide with ADR content using Title and Content layout."""
    layout = prs.slide_layouts[6]  # "Title and Content"
    slide = prs.slides.add_slide(layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title_text
    for run in title_shape.text_frame.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = COLOR_DARK

    # Find the body placeholder
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body = ph
            break
    if body is None:
        # Fallback: use any non-title placeholder
        for ph in slide.placeholders:
            if ph != title_shape:
                body = ph
                break

    if body is None:
        print(f"  WARNING: No body placeholder found for '{title_text}'")
        return slide

    # Clear default content
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True

    first = True
    for section_title, lines in sections:
        # Section header
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(10)

        run = p.add_run()
        run.text = section_title
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLOR_GREEN

        # Section content
        for line in lines:
            p = tf.add_paragraph()
            p.space_before = Pt(2)
            # Handle bullet prefix
            if line.startswith("- "):
                line = line[2:]
                p.level = 1
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            run.font.color.rgb = COLOR_DARK

    return slide


# ADR-001 slide
adr001_slide = add_adr_slide(prs, "ADR-001: HasFollowUp Trait Redesign", [
    ("Context", [
        "- HasFollowUp trait definieerde timeline scopes (scopeOverdue, scopeDueToday, etc.)",
        "- Scopes bevroegen direct follow_up_date en status kolommen",
        "- Deze kolommen bestaan alleen op follow_ups tabel; niet op Task of TeamMember",
        "- Runtime SQL errors gegarandeerd bij gebruik van de trait",
    ]),
    ("Alternatieven overwogen", [
        "- 1. Trait verwijderen; elk model definieert followUps() handmatig",
        "- 2. Scopes herschrijven naar whereHas-based queries",
    ]),
    ("Decision", [
        "- Optie 2: whereHas-based scopes die door de relatie filteren",
        "- Scope namen verduidelijkt (bijv. scopeWithOverdueFollowUps)",
        "- FollowUpStatus enum i.p.v. hardcoded strings",
        "- Trait nu daadwerkelijk gebruikt door Task en TeamMember",
    ]),
    ("Consequences", [
        "- Gewijzigd: HasFollowUp.php, Task.php, TeamMember.php",
        "- Geen database- of migratiewijzigingen nodig",
        "- Tests herschreven: testen nu whereHas scopes op Task als parent model",
    ]),
])

# ADR-013 slide
adr013_slide = add_adr_slide(prs, "ADR-013: Event Dispatch Gap & Recurring Tasks", [
    ("Context", [
        "- TaskStatusChanged event alleen gefired vanuit expliciete dispatch calls",
        "- AutoSaveController, bulkUpdate() en move() updaten status zonder event",
        "- Recurring tasks listener en CreateFollowUpOnWaiting zouden niet vuren",
    ]),
    ("Decision", [
        "- TaskObserver dispatcht TaskStatusChanged automatisch bij elke statuswijziging",
        "- Copy-on-complete: Done task -> nieuw Open task met volgende deadline",
        "- Recurrence velden op tasks tabel; geen apart model",
        "- Handmatige event dispatches bewust verwijderd",
    ]),
    ("Deviation from plan", [
        "- Plan voorzag niet in de event dispatch gap",
        "- TaskObserver toegevoegd als Phase 0 voor de recurrence implementatie",
    ]),
    ("Consequences", [
        "- Migratie: 5 kolommen op tasks tabel",
        "- Alle Task::update() met statuswijziging vuren nu automatisch het event",
        "- Lost ook pre-bestaand probleem op met CreateFollowUpOnWaiting listener",
    ]),
])

print("Step 2 done: created ADR slides")


# ============================================================
# STEP 3: Reorder slides
# ============================================================
# Current state: original slides 0-24, ADR-001 at 25, ADR-013 at 26
# Target: ADR-001 at position 11 (after case study slide 10)
#         ADR-013 at position 14 (after case study slide 13, which was old 12)

sld_id_lst = prs.slides._sldIdLst
items = list(sld_id_lst)

adr001_elem = items[-2]  # ADR-001 (added first)
adr013_elem = items[-1]  # ADR-013 (added second)

# Move ADR-001 to position 11
sld_id_lst.remove(adr001_elem)
sld_id_lst.insert(11, adr001_elem)

# Move ADR-013 to position 14 (after old slide 12, now at 13)
sld_id_lst.remove(adr013_elem)
sld_id_lst.insert(14, adr013_elem)

print("Step 3 done: reordered slides")
print(f"  Total slides: {len(prs.slides)}")


# ============================================================
# STEP 4: Set speaker notes (new indices, 0-26)
# ============================================================

# New slide order:
# 0:  Title
# 1:  Agenda
# 2:  DEEL 1 transition
# 3:  AI heeft geheugenverlies
# 4:  Context Rot in de praktijk
# 5:  De oplossing: externaliseer context
# 6:  DEEL 2 transition
# 7:  Laag 1: PRDs
# 8:  Laag 2: Het Plan
# 9:  Laag 3: TDD
# 10: Case Study ADR-001
# 11: ADR-001 Full Text (NEW)
# 12: Laag 4: ADRs
# 13: Case Study ADR-013
# 14: ADR-013 Full Text (NEW)
# 15: Feedback Loop
# 16: DEEL 3 transition
# 17: Vijf lessen
# 18: Hoe begin je?
# 19: Mithril in cijfers
# 20: DEEL 4 transition
# 21: Van handmatig naar georchestreerd
# 22: Zes rollen
# 23: Hooks
# 24: Context rot minimaliseren
# 25: De kern blijft
# 26: Discussie / Q&A

notes = {
    0: """Welkom allemaal bij deze workshop over TDD, ADR en AI-Assisted Development. Wat we vandaag gaan bespreken komt uit een echt project: Mithril, een Progressive Web App die ik in drie weken heb gebouwd met Claude Code als AI-assistent. Dat leverde 418 commits op, 2142 tests met 4602 assertions, en 28 architectuurbeslissingen die allemaal gedocumenteerd zijn. Dit is geen theoretisch verhaal; alles wat ik laat zien komt direct uit de praktijk.""",

    1: """De workshop bestaat uit vier delen. In het eerste deel kijken we naar het probleem: context rot. Waarom AI-assistants zonder structuur je codebase beschadigen. Dan in deel twee de oplossing: vier concrete lagen; PRD, Plan, TDD en ADR. Elk met voorbeelden uit Mithril. In deel drie gaan we het praktisch maken met de vijf belangrijkste lessen en hoe je er morgen mee kunt beginnen. En in deel vier laat ik ANVIL zien, een orchestrator die deze hele flow automatiseert.""",

    2: """Laten we beginnen bij het begin. Voordat we over oplossingen praten, moeten we het probleem goed begrijpen. Want als je niet weet waarom AI-assistants falen bij langere projecten, kun je ook niet effectief ingrijpen.""",

    3: """AI coding assistants werken binnen een context window; dat is een beperkt geheugen met een vaste maximale grootte. Elke nieuwe chat begint blanco; de AI weet dan niets meer van wat je eerder hebt besproken.

Maar het probleem gaat verder dan alleen sessiewisselingen. Ook binnen een lopend gesprek degradeert de context op drie manieren.

Ten eerste: context compressie. Wanneer een gesprek het window nadert, moet de AI oudere berichten samenvatten of zelfs laten vallen. Nuances, constraints en eerdere afspraken verdwijnen daarbij stilletjes.

Ten tweede: attention decay. Hoe meer tekst er in het window zit, hoe minder gewicht de AI geeft aan informatie die ver weg staat. Instructies die je twintig berichten geleden hebt gegeven worden effectief vergeten, ook al zijn ze technisch nog aanwezig.

En ten derde: hallucinaties. Wanneer de relevante context verdwenen of verdund is, vult de AI de gaten met plausibel klinkende maar feitelijk onjuiste informatie. Het ergste daaraan is dat de AI niet weet dat hij het niet meer weet. Hij genereert zelfverzekerd antwoorden op basis van incomplete context.

Om het concreet te maken: stel je voor dat je een collega hebt die scherp begint aan de dag, maar gedurende de dag steeds meer vergeet. En het ergste is dat hij niet doorheeft dat hij dingen vergeet.""",

    4: """Het resultaat van deze drie problemen noemen we context rot: een geleidelijk verval van de kwaliteit van AI-output naarmate een project vordert.

In de praktijk betekent dit dat architectuurbeslissingen verloren gaan tussen sessies, of verdund raken binnen een sessie. De AI maakt dezelfde fouten opnieuw omdat hij niet meer weet dat je die al hebt gecorrigeerd. Code wordt inconsistent omdat patronen vergeten worden. De vraag "waarom is dit zo?" is niet meer te beantwoorden; niet door jou en niet door de AI. En na een lange sessie kan de AI zelfverzekerd code voorstellen die lijnrecht ingaat tegen afspraken van eerder in datzelfde gesprek.""",

    5: """De oplossing is eigenlijk heel eenvoudig: de AI vergeet, maar bestanden in je repository niet. De sleutel is om context te externaliseren; om alles wat belangrijk is vast te leggen in bestanden die elke sessie overleven.

Dat doen we in vier lagen. PRDs beantwoorden wat we bouwen en waarom: feature-eisen, doelen en acceptatiecriteria. Plannen beantwoorden hoe we het bouwen: technische scope, fases en voortgang. TDD beantwoordt of het werkt: gedrag, contracten en regressiebeveiliging. En ADRs beantwoorden waarom we het zo doen: architectuurbeslissingen en hun onderbouwing.

Dit zijn de vier lagen die we nu stuk voor stuk gaan doorlopen.""",

    6: """Dan gaan we nu naar deel twee, waar we elke laag bekijken met concrete voorbeelden uit het Mithril project.""",

    7: """De eerste laag: PRDs, de opdracht. Alles begint met de vraag: wat gaan we bouwen en waarom?

Een PRD formaliseert dat. Het legt het doel vast; welk probleem lost dit op? De scope; wat is in en wat is out of scope? En de acceptatiecriteria; wanneer is het klaar?

Zonder PRD begint de AI, en soms ook de developer, te bouwen op basis van aannames. De PRD is het contract tussen het idee en de implementatie. Het voorkomt dat er een technisch plan wordt geschreven voor iets dat nog niet goed gedefinieerd is.

Belangrijk hierbij: een PRD hoeft niet zwaar te zijn. Zelfs een half A4 met een doel, een scope en vijf acceptatiecriteria bespaart je uren aan verkeerde implementatie.""",

    8: """De tweede laag: het plan. Met een PRD in de hand is de volgende stap: hoe gaan we dit bouwen?

Het plan vertaalt de criteria uit de PRD naar technische specificaties. Het deelt werk op in fases die in een enkele AI-sessie passen. Het maakt elke spec testbaar, wat directe input is voor TDD. En het houdt voortgang bij en legt vast wat out of scope is.

Het plan is ook een levend document. Tijdens de implementatie van recurring tasks in Mithril ontdekte TDD bijvoorbeeld een event dispatch gap die niet in het oorspronkelijke plan stond. Het plan werd bijgewerkt met een nieuwe Phase 0, en er werd een ADR geschreven. Dat voorbeeld komt straks terug als case study.""",

    9: """De derde laag: TDD, het veiligheidsnet. En ik zeg bewust veiligheidsnet; TDD is geen nice-to-have bij AI-assisted development. Het is het mechanisme dat voorkomt dat de AI je codebase sloopt.

Zonder TDD merk je niet dat de AI in sessie vijf iets heeft gebroken dat in sessie twee werkte. Met TDD worden regressies direct gevangen, en de test suite vormt een levende specificatie die elke sessie overleeft.

De cyclus is steeds dezelfde: Red, Green, Refactor. Eerst schrijf je een falende test die het gewenste gedrag beschrijft. Dan schrijf je de minimale code om die test te laten slagen. En tenslotte verbeter je de code met het vangnet van groene tests.""",

    10: """Laat me dit concreet maken met een echt voorbeeld uit Mithril. De HasFollowUp trait definieerde scopes die direct de kolommen follow_up_date en status bevroegen. Maar die kolommen bestaan alleen op de follow_ups tabel; niet op tasks of team_members, de modellen die de trait gebruiken. Dat betekent gegarandeerde SQL errors op runtime.

TDD ving dit. De test faalde niet omdat de implementatie fout was, maar omdat het hele design fout was. De scopes werden herschreven naar whereHas-based queries die correct door de relatie heen filteren. En de bug was gevonden voordat er ook maar een regel productcode was geschreven.

De belangrijkste les hier: de tests vonden de bug. Niet de AI, niet code review, niet QA. De falende test dwong een redesign af.""",

    11: """Dit is de daadwerkelijke ADR zoals die in de repository staat. Kijk naar de structuur: context, alternatieven, decision, consequences.

In de context beschrijven we het probleem dat TDD aan het licht bracht. Dan de alternatieven die we hebben overwogen: de trait helemaal verwijderen, of de scopes herschrijven. Bij de decision leggen we vast waarom we voor optie twee hebben gekozen. En tenslotte de consequences: welke bestanden zijn gewijzigd en wat verandert er voor toekomstige code.

Dit hele document is minder dan een pagina, maar het beantwoordt de vraag "waarom is dit zo?" voor elke toekomstige sessie. En dat is precies het punt: het kost vijf minuten om te schrijven, maar het bespaart uren aan herhaling en verwarring.""",

    12: """En daarmee zijn we bij de vierde laag: ADRs, het geheugen van je project. Tijdens het bouwen worden voortdurend beslissingen genomen die afwijken van het plan, of die niet-triviale trade-offs bevatten.

Een ADR documenteert drie dingen. De context: wat was het probleem? De decision: wat is er besloten en waarom? En de consequences: wat verandert er?

Zonder ADRs kan de AI in een volgende sessie een beslissing terugdraaien die je bewust hebt genomen, of een alternatief voorstellen dat je al hebt afgewezen.

De vuistregel is simpel: als iemand later zou kunnen vragen "waarom is dit zo?", schrijf dan een ADR.""",

    13: """Het tweede voorbeeld: de event dispatch gap. Bij het implementeren van recurring tasks bleek dat het TaskStatusChanged event alleen werd gefired vanuit expliciete dispatch calls in de code. Maar drie andere code paths; de AutoSaveController, bulkUpdate en de kanban move; updaten de status via een gewone update call zonder het event te dispatchen.

De oplossing was een TaskObserver die het event automatisch dispatcht bij elke statuswijziging, ongeacht welk code path de wijziging triggert. De handmatige dispatches zijn bewust verwijderd; er is nu een single source of truth.

Zonder deze ADR zou een volgende sessie niet weten dat er bewust een Observer is gekozen boven handmatige dispatch calls. Dat handmatige dispatches bewust zijn verwijderd. En welke alternatieven zijn afgewezen en waarom.""",

    14: """En hier zien jullie de volledige ADR-013 zoals die in het project staat. Let op het kopje "Deviation from plan"; dat is extra waardevol. Het legt expliciet vast waar de implementatie afweek van het oorspronkelijke plan, en waarom.

In dit geval had het plan niet voorzien dat er een event dispatch gap was in drie bestaande controllers. De Observer werd daarom als Phase 0 toegevoegd, nog voor de eigenlijke recurring tasks implementatie.

Dit soort informatie gaat verloren zodra een AI-sessie eindigt; tenzij je het vastlegt. En juist deze afwijkingen zijn het meest waardevol om te documenteren, want ze bevatten kennis die nergens anders te vinden is.""",

    15: """Nu we alle vier de lagen hebben gezien, laat me tonen hoe ze samenwerken in een feedback loop.

Het begint bij de PRD die definieert wat en waarom. Die wordt vertaald naar een plan met technische specs. Die specs worden tests via TDD. De tests ontdekken problemen. Die problemen triggeren ADRs. En die ADRs informeren de volgende sessies, waardoor de implementatie consistent blijft.

Concreet: de AI vergeet de vorige sessie, maar ADRs en het plan overleven sessies. De AI breekt bestaande code, maar tests vangen regressies direct. De AI gebruikt andere patronen, maar het plan en codebase scans dwingen patronen af. De AI voegt ongeplande features toe, maar het plan heeft een out-of-scope sectie. En de AI draait bewuste keuzes terug, maar ADRs documenteren het waarom.""",

    16: """Dan gaan we nu naar deel drie: toepassen en valkuilen. Hoe kun je dit in de praktijk gebruiken?""",

    17: """Dit zijn de vijf belangrijkste lessen uit drie weken bouwen met een AI-assistent.

Een: als het niet in een bestand staat, bestaat het niet. Elke AI-sessie begint leeg. Context die alleen in je hoofd zit, gaat gegarandeerd verloren.

Twee: maak specs testbaar. Als je iets niet in een test kunt uitdrukken, is de specificatie te vaag. Dit is waar plan en TDD samenkomen.

Drie: schrijf ADRs op het moment van de beslissing. Achteraf mis je de alternatieven die je hebt overwogen, en juist die alternatieven maken een ADR waardevol.

Vier: houd sessies klein. Context windows hebben limieten. Kleine fases in het plan houden sessies beheersbaar en voorkomen context compressie.

En vijf: jij blijft de architect. De AI is een krachtig hulpmiddel, maar review, denk mee, en stel vragen. Blind vertrouwen leidt tot blinde vlekken.""",

    18: """Het goede nieuws: je hebt geen speciale tooling nodig om hiermee te beginnen.

Start elke AI-sessie met context laden. Geef de AI bij de start de relevante ADRs, het plan, en laat hem de tests lezen.

Gebruik een CLAUDE.md of equivalent; een bestand in je repo dat de AI automatisch leest bij elke sessie. Hierin staan je conventies, actieve plannen, en verwijzingen naar ADRs.

Houd ADRs klein; een ADR per beslissing, niet een mega-document per feature.

En gebruik het plan als communicatiemiddel; niet alleen voor de AI, maar ook voor je team.""",

    19: """Om het concreet te maken: dit zijn de cijfers van Mithril. 418 commits, 2142 tests met 4602 assertions, 28 ADRs en 19 implementatieplannen. Allemaal in drie weken, met consistente architectuur ondanks tientallen sessiewisselingen. Dit laat zien dat de aanpak schaalbaar is; ook bij een project van deze omvang houdt de kwaliteit stand.""",

    20: """Dan het laatste deel: ANVIL, Automated Nexus for Verified Iterative Lifecycles. Van handmatig naar geautomatiseerd.""",

    21: """De vier lagen uit deel twee werken, maar ze vereisen discipline. Je moet zelf onthouden om een PRD te schrijven, een plan te maken, TDD te volgen, en ADRs bij te houden. In de praktijk slipt dat; vooral onder tijdsdruk.

ANVIL is een skill voor Claude Code die deze hele flow automatiseert. Het is in actieve ontwikkeling, gebouwd vanuit de lessen van het Mithril project.

ANVIL is een orchestrator die de volledige lifecycle aanstuurt via gespecialiseerde rollen. Van research door een Analyst, naar PRD door een Specifier, naar plan door een Architect, naar implementatie door een Builder met verplichte TDD, tot automatische ADRs bij afwijkingen.

Elke rol is een AI-agent met een specifieke opdracht en regels. De orchestrator stuurt ze aan, bewaakt de voortgang, en zorgt dat de gebruiker altijd de controle houdt.""",

    22: """ANVIL heeft zes gespecialiseerde rollen. De Analyst haalt requirements op uit externe bronnen zoals Jira. De Specifier schrijft en reviewt PRDs. De Architect maakt implementatieplannen met fases en specs. De Builder bouwt via TDD en kan parallel werken in agent teams. De Auditor doet OWASP security scans op de codebase. En de Scribe genereert documenten vanuit de opgeleverde resultaten.""",

    23: """Naast de rollen heeft ANVIL twee hooks die real-time afdwingen wat anders op discipline leunt.

De TDD Enforce hook waarschuwt direct wanneer je een bronbestand wijzigt zonder dat er testwijzigingen in de working tree staan. Je krijgt de melding op het moment van de edit, niet pas bij een commit. Dit dwingt de Red-Green volgorde af: eerst tests, dan implementatie.

De ADR Watch hook waarschuwt wanneer een bestand wordt gewijzigd dat buiten het actieve plan valt. Dit voorkomt onbewuste scope creep.

Beide hooks zijn informatief; ze blokkeren niet. En ze worden automatisch geinstalleerd. Het zijn mechanische checks die onafhankelijk van de AI draaien; geen herinneringen die de AI kan vergeten.""",

    24: """Hoe minimaliseert ANVIL context rot? Via drie mechanismen.

Ten eerste: state op disk, niet in het geheugen. ANVIL leidt de lifecycle state af uit bestanden. Niets hangt af van wat de AI onthoudt. Bij het einde van een sessie schrijft ANVIL een session handoff. De volgende sessie leest dat bestand en pikt het werk op waar het was gebleven.

Ten tweede: kleine, gefocuste agents. In plaats van een enkele AI-sessie die steeds langer wordt, splitst ANVIL het werk op. Een Builder-agent krijgt alleen de specs voor zijn fase, de relevante ADRs, en de projectconventies. Een kort, gefocust window betekent minder attention decay. Bijkomend voordeel: minder tokens per aanroep, lagere kosten, snellere respons.

En ten derde: mechanische bewaking vervangt geheugen. De hooks vangen fouten door context rot op het moment dat ze gebeuren, onafhankelijk van de AI.""",

    25: """Om af te sluiten: ANVIL automatiseert de flow, maar de principes zijn precies dezelfde als in deel twee. PRDs, plannen, tests en beslissingen; vastgelegd in bestanden die elke sessie overleven.

De tooling maakt het makkelijker om de discipline vol te houden, maar je kunt precies dezelfde aanpak toepassen met elk AI-hulpmiddel en een teksteditor.

ANVIL is op dit moment nog in actieve ontwikkeling. Ik zal de skill binnenkort beschikbaar stellen voor iedereen binnen Proud Nerds die hem wil uitproberen. Het is een Claude Code skill; je hebt een werkende Claude Code setup nodig.""",

    26: """Dan wil ik nu graag opengooien voor discussie en vragen. Een paar mogelijke gesprekspunten: hoe past dit in jullie huidige workflow? Welke onderdelen zijn direct toepasbaar, ook zonder AI? En wat denken jullie over de overhead versus de tijdsbesparing? Alle vragen zijn welkom.""",
}

for idx, notes_text in notes.items():
    slide = prs.slides[idx]
    if not slide.has_notes_slide:
        slide.notes_slide  # Access creates it
    slide.notes_slide.notes_text_frame.text = notes_text

print(f"Step 4 done: added notes to {len(notes)} slides")


# ============================================================
# STEP 5: Save
# ============================================================
prs.save(PPTX_PATH)
print(f"\nDone! Saved {PPTX_PATH} with {len(prs.slides)} slides.")
